"""Unified document converter using Docling."""

import os
from pathlib import Path
from typing import Dict, List
import zipfile
import tempfile
import shutil
import re
from loguru import logger
from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions, TableFormerMode
from PIL import Image
import io

from .base_converter import BaseConverter, DocumentContent
from .pptx_image_extractor import extract_images_from_pptx


class DoclingUnifiedConverter(BaseConverter):
    """Unified converter using Docling for all document types."""
    
    def __init__(self, config: Dict):
        """Initialize Docling converter.
        
        Args:
            config: Processing configuration dictionary
        """
        super().__init__(config)

        # Docling uses local GPU for layout analysis. Ollama runs on a separate host.
        import torch
        if torch.cuda.is_available():
            os.environ.setdefault('DOCLING_DEVICE', 'cuda')
            logger.info(f"Docling using GPU: {torch.cuda.get_device_name(0)}")

        # Page thresholds for mode selection
        self.large_pdf_threshold = config.get('docling', {}).get('large_pdf_page_threshold', 30)
        self.xlarge_pdf_threshold = config.get('docling', {}).get('xlarge_pdf_page_threshold', 80)

        # Default table mode from config
        table_mode_str = config.get('docling', {}).get('table_mode', 'fast').upper()
        table_mode = TableFormerMode.ACCURATE if table_mode_str == 'ACCURATE' else TableFormerMode.FAST

        # Standard converter with configured table mode
        pipeline_options = PdfPipelineOptions()
        pipeline_options.do_table_structure = True
        pipeline_options.table_structure_options.do_cell_matching = False
        pipeline_options.table_structure_options.mode = table_mode

        # FAST converter for large PDFs (30-80 pages)
        fast_pipeline = PdfPipelineOptions()
        fast_pipeline.do_table_structure = True
        fast_pipeline.table_structure_options.do_cell_matching = False
        fast_pipeline.table_structure_options.mode = TableFormerMode.FAST

        # Lightweight converter for very large PDFs (80+ pages) — no table structure at all
        lite_pipeline = PdfPipelineOptions()
        lite_pipeline.do_table_structure = False

        self.converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
            }
        )
        self.fast_converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=fast_pipeline)
            }
        )
        self.lite_converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=lite_pipeline)
            }
        )

        logger.info(f"Docling initialized (table mode: {table_mode_str}, large: >{self.large_pdf_threshold}p=FAST, xlarge: >{self.xlarge_pdf_threshold}p=LITE)")
    
    def convert(self, file_path: Path) -> DocumentContent:
        """Convert document using Docling.
        
        Args:
            file_path: Path to document file
            
        Returns:
            DocumentContent object with extracted content
        """
        logger.info(f"Converting with Docling: {file_path}")
        content = DocumentContent()
        
        try:
            # Log file details for debugging
            logger.debug(f"File path: {file_path}")
            logger.debug(f"File exists: {file_path.exists()}")
            logger.debug(f"File size: {file_path.stat().st_size if file_path.exists() else 'N/A'}")
            
            # Select converter based on PDF page count
            converter = self.converter
            if file_path.suffix.lower() == '.pdf':
                try:
                    import pypdfium2 as pdfium
                    pdf = pdfium.PdfDocument(str(file_path))
                    page_count = len(pdf)
                    pdf.close()
                    if page_count > self.xlarge_pdf_threshold:
                        logger.warning(f"Very large PDF ({page_count} pages > {self.xlarge_pdf_threshold}): using LITE mode (no table structure)")
                        converter = self.lite_converter
                    elif page_count > self.large_pdf_threshold:
                        logger.warning(f"Large PDF ({page_count} pages > {self.large_pdf_threshold}): using FAST table mode")
                        converter = self.fast_converter
                    else:
                        logger.info(f"PDF has {page_count} pages, using default table mode")
                except Exception as page_err:
                    logger.debug(f"Could not count pages: {page_err}")

            # Convert document using Docling
            result = converter.convert(str(file_path))
            
            # Extract metadata
            if hasattr(result.document, 'metadata'):
                metadata = result.document.metadata
                content.title = getattr(metadata, 'title', file_path.stem) or file_path.stem
                
                if hasattr(metadata, 'author'):
                    content.set_metadata('author', metadata.author)
                if hasattr(metadata, 'subject'):
                    content.set_metadata('subject', metadata.subject)
                if hasattr(metadata, 'creator'):
                    content.set_metadata('creator', metadata.creator)
            else:
                content.title = file_path.stem
            
            # Get markdown output from Docling
            markdown_text = result.document.export_to_markdown()
            content.add_text(markdown_text)
            
            # Extract images with Docling
            image_count = 0
            if self.extract_images and hasattr(result.document, 'pictures'):
                for picture in result.document.pictures:
                    try:
                        # Get image data - try multiple extraction methods
                        pil_image = None
                        
                        # Method 1: Direct image attribute (newer Docling API)
                        if hasattr(picture, 'image') and picture.image:
                            if isinstance(picture.image, Image.Image):
                                pil_image = picture.image
                            else:
                                logger.debug(f"picture.image is not a PIL Image: {type(picture.image)}")
                        
                        # Method 2: get_image() method (older API)
                        if pil_image is None and hasattr(picture, 'get_image'):
                            try:
                                image_ref = picture.get_image(result.document)
                                # Try various attributes of image_ref
                                if hasattr(image_ref, 'pil_image'):
                                    pil_image = image_ref.pil_image
                                elif hasattr(image_ref, 'data'):
                                    pil_image = Image.open(io.BytesIO(image_ref.data))
                                elif hasattr(image_ref, 'image'):
                                    pil_image = image_ref.image
                            except Exception as img_err:
                                logger.debug(f"get_image() failed: {img_err}")
                        
                        # Method 3: Try to export picture directly
                        if pil_image is None and hasattr(picture, 'export_to_pil'):
                            try:
                                pil_image = picture.export_to_pil()
                            except:
                                pass
                        
                        if pil_image is None:
                            logger.debug(f"Could not extract image {image_count}, skipping")
                            continue
                        
                        # Get caption/alt text if available
                        caption = getattr(picture, 'caption', '')
                        
                        # Always try OCR for presentation slides (they often contain text as images)
                        ocr_text = ""
                        if self.ocr_enabled:
                            ocr_text = self.perform_ocr(pil_image)
                            if ocr_text:
                                logger.info(f"OCR extracted {len(ocr_text)} chars from image {image_count}")
                        
                        # Save image to temp directory for AI vision analysis
                        img_format = pil_image.format if pil_image.format else 'PNG'
                        img_filename = f"{file_path.stem}_image_{image_count}.{img_format.lower()}"
                        img_path = self.temp_dir / img_filename
                        pil_image.save(img_path, format=img_format)
                        
                        image_info = {
                            'index': image_count,
                            'caption': caption,
                            'has_ocr_text': bool(ocr_text),
                            'ocr_text': ocr_text,
                            'image': pil_image,
                            'path': str(img_path),
                            'filename': img_filename,
                            'format': img_format,
                            'size': pil_image.size
                        }
                        content.images.append(image_info)
                        image_count += 1
                        
                    except Exception as e:
                        logger.warning(f"Failed to extract image {image_count}: {e}")
                
                logger.info(f"Extracted {image_count} images with Docling")
            
            # Fallback: For PPTX files, try python-pptx if no images extracted
            if image_count == 0 and file_path.suffix.lower() in ['.pptx', '.ppt']:
                logger.info("Attempting PPTX image extraction with python-pptx...")
                # Use temporary directory from config
                temp_dir = Path(self.config.get('temp_dir', '/tmp/docling'))
                temp_dir.mkdir(parents=True, exist_ok=True)
                pptx_images = extract_images_from_pptx(file_path, temp_dir)
                if pptx_images:
                    # Perform OCR on extracted images
                    for img_info in pptx_images:
                        if self.ocr_enabled and img_info.get('image'):
                            ocr_text = self.perform_ocr(img_info['image'])
                            if ocr_text:
                                img_info['ocr_text'] = ocr_text
                                img_info['has_ocr_text'] = True
                                logger.info(f"OCR extracted {len(ocr_text)} chars from {img_info['filename']}")
                    
                    content.images.extend(pptx_images)
                    logger.info(f"✓ Extracted {len(pptx_images)} images from PPTX with fallback method")
            
            # Get page count if available
            if hasattr(result.document, 'pages'):
                content.set_metadata('page_count', len(result.document.pages))
            
            logger.info(f"Docling conversion complete for: {file_path.name}")
            
        except Exception as e:
            import traceback
            error_traceback = traceback.format_exc()
            logger.error(f"Error converting with Docling: {e}")
            logger.error(f"Full traceback:\n{error_traceback}")
            
            # For PPTX files, always try auto-repair first (Docling may fail due to XML issues)
            # The XMLSyntaxError happens inside Docling before it raises ConversionError to us
            if file_path.suffix.lower() in ['.pptx', '.ppt']:
                logger.warning(f"Docling failed on PPTX, attempting auto-repair for potential XML issues...")
                try:
                    repaired_path = self._repair_pptx_xml(file_path)
                    if repaired_path:
                        logger.info(f"Successfully repaired PPTX, retrying conversion...")
                        # Retry conversion with repaired file
                        try:
                            result = self.converter.convert(str(repaired_path))
                            
                            # Extract metadata (same as above)
                            if hasattr(result.document, 'metadata'):
                                metadata = result.document.metadata
                                content.title = getattr(metadata, 'title', file_path.stem) or file_path.stem
                                
                                if hasattr(metadata, 'author'):
                                    content.set_metadata('author', metadata.author)
                                if hasattr(metadata, 'subject'):
                                    content.set_metadata('subject', metadata.subject)
                                if hasattr(metadata, 'creator'):
                                    content.set_metadata('creator', metadata.creator)
                            else:
                                content.title = file_path.stem
                            
                            # Get markdown
                            markdown_text = result.document.export_to_markdown()
                            content.add_text(markdown_text)
                            
                            # Extract images via pptx_image_extractor
                            temp_dir = Path(self.config.get('temp_dir', '/tmp/docling'))
                            temp_dir.mkdir(parents=True, exist_ok=True)
                            content.images = extract_images_from_pptx(file_path, temp_dir)
                            
                            # Get page count
                            if hasattr(result.document, 'pages'):
                                content.set_metadata('page_count', len(result.document.pages))
                            
                            logger.success(f"✓ Converted repaired PPTX successfully!")
                            return content
                            
                        except Exception as retry_err:
                            logger.error(f"Repaired file also failed: {retry_err}")
                            # Continue to fallback below
                        finally:
                            # Clean up repaired file
                            if repaired_path.exists():
                                repaired_path.unlink()
                except Exception as repair_err:
                    logger.error(f"Auto-repair failed: {repair_err}")
            
            # Try python-pptx fallback for PPTX files that Docling can't handle
            if file_path.suffix.lower() in ['.pptx', '.ppt']:
                logger.info(f"Attempting python-pptx fallback for: {file_path.name}")
                try:
                    from pptx import Presentation
                    fallback_content = DocumentContent()
                    fallback_content.title = file_path.stem
                    
                    prs = Presentation(file_path)
                    
                    # Extract text from all slides
                    slide_texts = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_text = f"\n## Slide {slide_num}\n\n"
                        
                        # Extract text from shapes
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_text += f"{shape.text}\n\n"
                        
                        # Extract speaker notes
                        if slide.has_notes_slide:
                            notes_slide = slide.notes_slide
                            if notes_slide.notes_text_frame:
                                speaker_notes = notes_slide.notes_text_frame.text.strip()
                                if speaker_notes:
                                    slide_text += f"**Speaker Notes:**\n> {speaker_notes}\n\n"
                        
                        slide_texts.append(slide_text)
                    
                    fallback_content.add_text("\n".join(slide_texts))
                    
                    # Extract images using pptx_image_extractor
                    temp_dir = Path(self.config.get('temp_dir', '/tmp/docling'))
                    temp_dir.mkdir(parents=True, exist_ok=True)
                    images = extract_images_from_pptx(file_path, temp_dir)
                    fallback_content.images = images
                    
                    logger.success(f"✓ Python-pptx fallback successful: {len(prs.slides)} slides, {len(images)} images")
                    return fallback_content
                    
                except Exception as fallback_err:
                    logger.error(f"Python-pptx fallback also failed: {fallback_err}")
                    content.add_text(f"Error converting document: {str(e)}\nFallback also failed: {str(fallback_err)}")
                    return content
            
            # For non-PPTX, return error content
            content.add_text(f"Error converting document: {str(e)}")
        
        return content
    
    def _repair_pptx_xml(self, file_path: Path) -> Path:
        """Repair corrupted XML in PPTX files.
        
        Args:
            file_path: Path to corrupted PPTX
            
        Returns:
            Path to repaired PPTX file
        """
        try:
            # Create temporary directory for repair
            temp_dir = tempfile.mkdtemp()
            extract_dir = Path(temp_dir) / "extracted"
            extract_dir.mkdir()
            
            # Extract PPTX (it's a ZIP file)
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)
            
            logger.debug(f"Extracted PPTX to {extract_dir}")
            
            # Find and fix XML files with namespace errors
            fixed_count = 0
            for xml_file in extract_dir.rglob("*.xml"):
                try:
                    with open(xml_file, 'r', encoding='utf-8') as f:
                        content = f.read()
                    
                    # Fix common XML namespace issues
                    original_content = content
                    
                    # Fix trailing backslashes in namespace URIs
                    # Pattern: xmlns:prefix="http://...\"
                    content = re.sub(r'(xmlns:\w+="[^"]+)\\(")', r'\1\2', content)
                    
                    # Fix other common namespace issues
                    content = re.sub(r'(xmlns:\w+="[^"]+)\s+\\', r'\1', content)
                    
                    if content != original_content:
                        with open(xml_file, 'w', encoding='utf-8') as f:
                            f.write(content)
                        fixed_count += 1
                        logger.debug(f"Fixed XML in: {xml_file.name}")
                
                except Exception as file_err:
                    logger.debug(f"Could not process {xml_file.name}: {file_err}")
                    continue
            
            if fixed_count == 0:
                logger.warning("No XML fixes applied")
                shutil.rmtree(temp_dir)
                return None
            
            logger.info(f"Fixed {fixed_count} XML file(s)")
            
            # Create repaired PPTX
            repaired_path = Path(temp_dir) / f"repaired_{file_path.name}"
            with zipfile.ZipFile(repaired_path, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
                for file in extract_dir.rglob("*"):
                    if file.is_file():
                        arcname = file.relative_to(extract_dir)
                        zip_ref.write(file, arcname)
            
            logger.success(f"Created repaired PPTX: {repaired_path}")
            return repaired_path
            
        except Exception as e:
            logger.error(f"Failed to repair PPTX: {e}")
            if 'temp_dir' in locals():
                shutil.rmtree(temp_dir, ignore_errors=True)
            return None
    
    def supports_format(self, file_path: Path) -> bool:
        """Check if file format is supported.
        
        Args:
            file_path: Path to file
            
        Returns:
            True if supported, False otherwise
        """
        ext = file_path.suffix.lower()
        supported = ['.pdf', '.docx', '.pptx', '.html', '.htm', '.png', '.jpg', '.jpeg']
        return ext in supported


class DoclingPDFConverter(DoclingUnifiedConverter):
    """Docling-based PDF converter (for compatibility)."""
    pass


class DoclingDocxConverter(DoclingUnifiedConverter):
    """Docling-based DOCX converter (for compatibility)."""
    pass


class DoclingPptxConverter(DoclingUnifiedConverter):
    """Docling-based PPTX converter (for compatibility)."""
    pass

