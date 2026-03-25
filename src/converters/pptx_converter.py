"""PPTX presentation converter."""

from pathlib import Path
from typing import Dict
from loguru import logger
from pptx import Presentation
from PIL import Image
import io

from .base_converter import BaseConverter, DocumentContent


class PptxConverter(BaseConverter):
    """Converts PPTX presentations to markdown."""
    
    def __init__(self, config: Dict):
        """Initialize PPTX converter.
        
        Args:
            config: Processing configuration dictionary
        """
        super().__init__(config)
    
    def convert(self, file_path: Path) -> DocumentContent:
        """Convert PPTX to markdown content.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            DocumentContent object with extracted content
        """
        logger.info(f"Converting PPTX: {file_path}")
        content = DocumentContent()
        
        try:
            # Open presentation
            prs = Presentation(file_path)
            
            # Extract metadata
            core_props = prs.core_properties
            content.title = core_props.title or file_path.stem
            content.set_metadata('author', core_props.author or '')
            content.set_metadata('subject', core_props.subject or '')
            content.set_metadata('slide_count', len(prs.slides))
            
            # Process each slide
            image_count = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                logger.debug(f"Processing slide {slide_num}/{len(prs.slides)}")
                
                # Add slide header
                content.add_text(f"## Slide {slide_num}")
                
                # Extract text from shapes
                slide_text = []
                for shape in slide.shapes:
                    # Text content
                    if hasattr(shape, "text") and shape.text.strip():
                        # Check if it's a title
                        if shape == slide.shapes.title:
                            slide_text.insert(0, f"**{shape.text.strip()}**")
                        else:
                            slide_text.append(shape.text.strip())
                    
                    # Extract images from shapes
                    if self.extract_images and hasattr(shape, "image"):
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            pil_image = Image.open(io.BytesIO(image_bytes))
                            
                            # Try OCR
                            ocr_text = self.perform_ocr(pil_image)
                            
                            image_info = {
                                'slide': slide_num,
                                'index': image_count,
                                'has_ocr_text': bool(ocr_text),
                                'ocr_text': ocr_text,
                                'image': pil_image,
                                'format': image.content_type.split('/')[-1]
                            }
                            content.images.append(image_info)
                            image_count += 1
                            
                        except Exception as e:
                            logger.warning(f"Failed to extract image from slide {slide_num}: {e}")
                
                # Add slide text
                if slide_text:
                    content.add_text("\n\n".join(slide_text))
                
                # Check for notes
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        content.add_text(f"\n**Notes:** {notes_text.strip()}")
                
                # Add separator between slides
                if slide_num < len(prs.slides):
                    content.add_text("---")
            
            logger.info(f"PPTX conversion complete. Extracted {image_count} images")
            
        except Exception as e:
            logger.error(f"Error converting PPTX: {e}")
            content.add_text(f"Error converting PPTX: {str(e)}")
        
        return content

