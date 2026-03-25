"""PowerPoint image extraction helper using python-pptx."""

from pathlib import Path
from typing import List, Dict
from PIL import Image
import io
from loguru import logger

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available, PPTX image extraction limited")


def extract_images_from_pptx(file_path: Path, temp_dir: Path) -> List[Dict]:
    """Extract images AND speaker notes from PPTX file using python-pptx.
    
    This is a fallback when Docling can't extract images.
    
    Args:
        file_path: Path to PPTX file
        temp_dir: Temporary directory for saving images
        
    Returns:
        List of image info dictionaries with speaker notes
    """
    if not PPTX_AVAILABLE:
        logger.warning("python-pptx not available")
        return []
    
    images = []
    
    try:
        prs = Presentation(str(file_path))
        image_count = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Extract speaker notes for this slide
            speaker_notes = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    speaker_notes = notes_slide.notes_text_frame.text.strip()
                    if speaker_notes:
                        logger.info(f"Extracted {len(speaker_notes)} chars of speaker notes from slide {slide_num}")
            
            for shape in slide.shapes:
                # Check if shape has an image
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        
                        # Get image bytes
                        image_bytes = image.blob
                        
                        # Convert to PIL Image
                        pil_image = Image.open(io.BytesIO(image_bytes))
                        
                        # Determine format
                        img_format = pil_image.format if pil_image.format else 'PNG'
                        
                        # Save image
                        img_filename = f"{file_path.stem}_slide{slide_num}_img{image_count}.{img_format.lower()}"
                        img_path = temp_dir / img_filename
                        pil_image.save(img_path, format=img_format)
                        
                        # Build caption with speaker notes
                        caption = f"Slide {slide_num}"
                        if speaker_notes:
                            caption += f" - Speaker Notes: {speaker_notes[:100]}..."
                        
                        images.append({
                            'index': image_count,
                            'slide_number': slide_num,
                            'caption': caption,
                            'speaker_notes': speaker_notes,
                            'has_ocr_text': False,
                            'ocr_text': '',
                            'image': pil_image,
                            'path': str(img_path),
                            'filename': img_filename,
                            'format': img_format,
                            'size': pil_image.size
                        })
                        
                        image_count += 1
                        logger.info(f"Extracted image from slide {slide_num}")
                        
                    except Exception as e:
                        logger.debug(f"Failed to extract image from slide {slide_num}: {e}")
        
        logger.info(f"python-pptx extracted {len(images)} images")
        return images
        
    except Exception as e:
        logger.error(f"Failed to extract images with python-pptx: {e}")
        return []

